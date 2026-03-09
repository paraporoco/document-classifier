from pathlib import Path
from docx import Document
import openpyxl
from pptx import Presentation

class ExtractionError(Exception):
    pass

def _docx(path):
    doc = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())
    return "\n".join(parts)

def _xlsx(path):
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = list(wb.sheetnames)
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    v = str(cell).strip()
                    if v and v != "None":
                        parts.append(v)
    wb.close()
    return "\n".join(parts)

def _pptx(path):
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)
    return "\n".join(parts)

def extract(path: Path) -> str:
    try:
        suffix = path.suffix.lower()
        if suffix == ".docx": return _docx(path)
        if suffix == ".xlsx": return _xlsx(path)
        if suffix == ".pptx": return _pptx(path)
        raise ExtractionError(f"Unsupported: {suffix}")
    except ExtractionError:
        raise
    except Exception as e:
        raise ExtractionError(str(e)) from e
