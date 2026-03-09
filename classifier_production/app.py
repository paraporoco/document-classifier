"""
classify — document sensitivity classifier
Usage: classify.exe <folder> [--output log.csv] [--recursive] [--review-only]
"""

import argparse, csv, os, sys
from datetime import datetime, timezone
from pathlib import Path
from collections import Counter

# ── Extractors ────────────────────────────────────────────────────────────────

def _read_text(path):
    try:    return path.read_text(encoding="utf-8"), "ok", None
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1"), "ok", None

def _read_pdf_pdfminer(path):
    from io import StringIO
    from pdfminer.high_level import extract_text_to_fp
    from pdfminer.layout import LAParams
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.pdfparser import PDFSyntaxError
    buf = StringIO()
    with open(path, "rb") as f:
        extract_text_to_fp(f, buf, laparams=LAParams(), output_type="text", codec="utf-8")
    return buf.getvalue().strip()

def _read_pdf_pypdf(path):
    import pypdf
    reader = pypdf.PdfReader(str(path))
    parts = []
    for page in reader.pages:
        t = page.extract_text()
        if t: parts.append(t)
    return "\n".join(parts).strip()

_tesseract_cmd = None   # set at startup from --tesseract arg

TESSERACT_DEFAULT = (
    r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if sys.platform == "win32"
    else "/usr/bin/tesseract"
)

def _read_pdf_ocr(path):
    """Render PDF pages to images via pymupdf and OCR with Tesseract."""
    import pytesseract
    import fitz  # pymupdf
    pytesseract.pytesseract.tesseract_cmd = _tesseract_cmd
    doc = fitz.open(str(path))
    parts = []
    for page in doc:
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("png")
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(img_bytes))
        text = pytesseract.image_to_string(img, lang="fra+eng").strip()
        if text:
            parts.append(text)
    doc.close()
    return "\n".join(parts).strip()

def _read_pdf(path):
    try:
        from pdfminer.pdfdocument import PDFPasswordIncorrect
        from pdfminer.pdfparser import PDFSyntaxError
        try:
            text = _read_pdf_pdfminer(path)
        except PDFPasswordIncorrect:
            return None, "encrypted", "Password-protected"
        except Exception:
            text = ""
        if len(text) < 20:
            # fallback to pypdf
            try:
                text = _read_pdf_pypdf(path)
            except Exception:
                text = ""   # pypdf also failed — treat as no text layer
        if len(text) < 20:
            # Last resort: OCR via Tesseract
            if _tesseract_cmd:
                try:
                    text = _read_pdf_ocr(path)
                except FileNotFoundError:
                    return None, "no_text_layer", (
                        f"Tesseract not found at: {_tesseract_cmd}  "
                        "Install from https://github.com/UB-Mannheim/tesseract/wiki "
                        "or use --tesseract to specify the path."
                    )
                except Exception as e:
                    return None, "no_text_layer", f"OCR failed: {type(e).__name__}: {e}"
            if len(text) < 20:
                if _tesseract_cmd:
                    detail = "OCR ran but returned no text (pure image with no recognisable characters?)"
                else:
                    detail = ("No extractable text. Install Tesseract and add --tesseract to enable OCR.  "
                              "Download: https://github.com/UB-Mannheim/tesseract/wiki")
                return None, "no_text_layer", detail
        return text, "ok", None
    except Exception as e:
        return None, "extraction_failed", f"{type(e).__name__}: {e}"

def _read_html(path):
    from html.parser import HTMLParser
    class _Strip(HTMLParser):
        def __init__(self):
            super().__init__()
            self._parts = []
            self._skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"): self._skip = True
        def handle_endtag(self, tag):
            if tag in ("script", "style"): self._skip = False
        def handle_data(self, data):
            if not self._skip:
                t = data.strip()
                if t: self._parts.append(t)
        def text(self): return "\n".join(self._parts)
    try:    raw = path.read_text(encoding="utf-8")
    except UnicodeDecodeError: raw = path.read_text(encoding="latin-1")
    p = _Strip(); p.feed(raw)
    return p.text(), "ok", None

def _read_docx(path):
    from docx import Document
    doc = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip(): parts.append(cell.text.strip())
    return "\n".join(parts), "ok", None

def _read_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = list(wb.sheetnames)
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    v = str(cell).strip()
                    if v and v != "None": parts.append(v)
    wb.close()
    return "\n".join(parts), "ok", None

def _read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip(): parts.append(para.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes: parts.append(notes)
    return "\n".join(parts), "ok", None

def extract(path: Path):
    if path.stat().st_size > 10 * 1024 * 1024:
        return None, "oversized", f"{path.stat().st_size/1024/1024:.1f} MB > 10 MB limit"
    s = path.suffix.lower()
    try:
        if s in (".txt", ".md"):   return _read_text(path)
        if s == ".pdf":            return _read_pdf(path)
        if s in (".html", ".htm"): return _read_html(path)
        if s == ".docx":           return _read_docx(path)
        if s == ".xlsx":           return _read_xlsx(path)
        if s == ".pptx":           return _read_pptx(path)
        return None, "unsupported_format", f"'{s}' not supported"
    except Exception as e:
        return None, "extraction_failed", str(e)

# ── Model ─────────────────────────────────────────────────────────────────────

import pickle
import numpy as np

CLASSES          = ["PUBLIC","FOUO","CONFIDENTIAL","PERSONAL_CONFIDENTIAL","HIGHLY_CONFIDENTIAL","PERSONAL_HIGHLY_CONFIDENTIAL"]
REVIEW_THRESHOLD = 0.80
POLICY_VERSION   = "1.1"
_model = _vec = _T = None

def _artefact(name):
    base = Path(getattr(sys, "_MEIPASS", Path(__file__).parent))
    return base / "artefacts" / name

def load_model():
    global _model, _vec, _T
    with open(_artefact("model_ssl.pkl"),   "rb") as f: _model = pickle.load(f)
    with open(_artefact("vectorizer.pkl"),  "rb") as f: _vec   = pickle.load(f)
    with open(_artefact("temperature.pkl"), "rb") as f: _T     = pickle.load(f)

def classify(text: str) -> dict:
    X = _vec.transform([text]).toarray().astype(np.float32)
    proba  = _model.predict_proba(X)
    logits = np.log(proba + 1e-12)
    scaled = logits / _T
    scaled -= scaled.max(axis=1, keepdims=True)
    exp    = np.exp(scaled)
    cal    = exp / exp.sum(axis=1, keepdims=True)
    idx    = int(cal.argmax(axis=1)[0])
    conf   = float(cal.max(axis=1)[0])
    names  = np.array(_vec.get_feature_names_out())
    w      = np.abs(_model.coefs_[0]).sum(axis=1)
    top    = (X[0] * w).argsort()[-5:][::-1]
    feats  = [f"{names[i]} ({X[0][i]*w[i]:.3f})" for i in top if X[0][i]*w[i] > 0]
    return {
        "classification" : CLASSES[idx],
        "confidence"     : round(conf, 4),
        "review_flag"    : conf < REVIEW_THRESHOLD,
        "reasoning"      : ("Top signals: " + ", ".join(feats)) if feats else "No strong signals.",
        "policy_version" : POLICY_VERSION,
    }

# ── Output ────────────────────────────────────────────────────────────────────

FIELDS = ["timestamp","filepath","filename","extension","size_bytes",
          "classification","confidence","review_flag","reasoning",
          "policy_version","extraction_status","extraction_detail"]

def _ts():
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

def _walk(folder, recursive):
    if recursive:
        for root, _, files in os.walk(folder):
            for f in sorted(files): yield Path(root) / f
    else:
        for f in sorted(folder.iterdir()):
            if f.is_file(): yield f

def _summary(output, rows):
    p = output.with_name(output.stem + "_summary.txt")
    classified = [r for r in rows if r["classification"]]
    errors     = [r for r in rows if not r["classification"]]
    reviews    = [r for r in rows if str(r.get("review_flag","")).lower() == "true"]
    counts     = Counter(r["classification"] for r in classified)
    lines = [
        "CLASSIFICATION RUN SUMMARY", "=" * 40,
        f"Run        : {_ts()}",
        f"Total      : {len(rows)}",
        f"Classified : {len(classified)}",
        f"Errors     : {len(errors)}",
        f"Flagged    : {len(reviews)}", "",
        "By level:",
    ]
    for cls in CLASSES:
        n = counts.get(cls, 0)
        lines.append(f"  {cls:<35} {'█'*n} {n}")
    if errors:
        lines += ["", "Errors:"]
        for s, n in Counter(r["extraction_status"] for r in errors).items():
            lines.append(f"  {s:<30} {n}")
    p.write_text("\n".join(lines), encoding="utf-8")
    print(f"  Summary  → {p}")

def _write_slim(output: Path, rows: list) -> None:
    slim_path = output.with_name(output.stem + "_slim.csv")
    with open(slim_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["filename", "classification"])
        writer.writeheader()
        writer.writerows({"filename": r["filename"], "classification": r["classification"]} for r in rows)
    print(f"  Slim log → {slim_path}")

# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(prog="classify",
             description="Classify documents by data sensitivity level.")
    ap.add_argument("input",            type=Path, help="File or folder to classify")
    ap.add_argument("--output",         type=Path, default=Path("classification_log.csv"))
    ap.add_argument("--recursive",      action="store_true")
    ap.add_argument("--review-only",    action="store_true")
    ap.add_argument("--min-confidence", type=float, default=0.0)
    ap.add_argument("--tesseract",      type=str,   default=None,
                    help="Path to tesseract.exe (default: standard install location)")
    args = ap.parse_args()

    if not args.input.exists():
        sys.exit(f"ERROR: not found: {args.input}")

    # Configure Tesseract path
    global _tesseract_cmd
    tess = args.tesseract or TESSERACT_DEFAULT
    if Path(tess).exists():
        _tesseract_cmd = tess
        print(f"Tesseract: {tess}")
    else:
        _tesseract_cmd = None
        print(f"Tesseract: not found at {tess} — OCR disabled for image PDFs")
        print(f"           Install from https://github.com/UB-Mannheim/tesseract/wiki")
        if args.tesseract:
            print(f"           or check the path passed to --tesseract")

    print("Loading model…")
    load_model()
    print("  Ready.")

    if args.input.is_file():
        files = [args.input]
        print(f"Classifying: {args.input}\n")
    else:
        files = list(_walk(args.input, args.recursive))
        print(f"Found {len(files)} file(s) in {args.input}\n")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    rows = []
    w = len(str(len(files)))

    with open(args.output, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=FIELDS)
        writer.writeheader()

        for i, path in enumerate(files, 1):
            label = f"  [{i:>{w}}/{len(files)}] {path.name:<45}"
            text, status, detail = extract(path)

            if status != "ok":
                row = {"timestamp": _ts(), "filepath": str(path), "filename": path.name,
                       "extension": path.suffix.lower(), "size_bytes": path.stat().st_size if path.exists() else "",
                       "classification": "", "confidence": "", "review_flag": "",
                       "reasoning": "", "policy_version": POLICY_VERSION,
                       "extraction_status": status, "extraction_detail": detail or ""}
                print(f"{label} ✗ {status}")
                if detail and status not in ("unsupported_format",):
                    print(f"  {'':>{len(label)-2}}  ↳ {detail}")
            else:
                try:
                    result = classify(text)
                except Exception as e:
                    row = {"timestamp": _ts(), "filepath": str(path), "filename": path.name,
                           "extension": path.suffix.lower(), "size_bytes": path.stat().st_size,
                           "classification": "", "confidence": "", "review_flag": "",
                           "reasoning": "", "policy_version": POLICY_VERSION,
                           "extraction_status": "classification_failed", "extraction_detail": str(e)}
                    print(f"{label} ✗ classification_failed")
                    rows.append(row); writer.writerow(row); continue

                if args.min_confidence > 0 and result["confidence"] < args.min_confidence:
                    print(f"{label} – skipped (conf={result['confidence']:.3f})"); continue
                if args.review_only and not result["review_flag"]:
                    print(f"{label} – skipped"); continue

                row = {"timestamp": _ts(), "filepath": str(path), "filename": path.name,
                       "extension": path.suffix.lower(), "size_bytes": path.stat().st_size,
                       "classification": result["classification"],
                       "confidence": result["confidence"],
                       "review_flag": result["review_flag"],
                       "reasoning": result["reasoning"],
                       "policy_version": result["policy_version"],
                       "extraction_status": "ok", "extraction_detail": ""}
                flag = " ⚑" if result["review_flag"] else ""
                print(f"{label} ✓ {result['classification']:<35} {result['confidence']:.3f}{flag}")

            rows.append(row)
            writer.writerow(row)

    print(f"\nLog → {args.output}")
    _summary(args.output, rows)
    _write_slim(args.output, rows)

if __name__ == "__main__":
    main()
